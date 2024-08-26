from pptx import Presentation
from pptx.shapes.picture import Picture
from PIL import Image, ImageStat
import pytesseract

import os, io
from pptx import Presentation
import streamlit as st

# def is_image_black(image):
#     # 检查图像是否是全黑
#     stat = ImageStat.Stat(image)
#     if sum(stat.extrema[0]) == 0:
#         return True
#     return False
# 假设你有一个类似的函数来提取图片和文字
def extract_text_from_slide(slide):
    """提取单张幻灯片的文本内容"""
    slide_text = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                slide_text.append(paragraph.text)
    return "\n".join(slide_text)

def extract_tables_from_slide(slide):
    """提取单张幻灯片的表格数据"""
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)
    return tables
def extract_images_from_slide(slide, slide_number,pptx_name):
    """提取单张幻灯片的图片"""
    images_with_text = []
    image_id = 0
    for shape_number, shape in enumerate(slide.shapes, start=1):
        text = ""
        if shape.shape_type == 13:  # Shape 类型 13 是图片
            image = shape.image
            image_bytes = image.blob
            image_ext = image.ext
            # 将图像加载到 PIL 图像对象
            image = Image.open(io.BytesIO(image_bytes))

            # if is_image_black(image):
            #     continue  # 跳过全黑图像
                # 将图片存储为文件
            image_filename = f"{pptx_name}_image_{slide_number}_{image_id + 1}.{image_ext}"
            image_id = image_id + 1
            with open(image_filename, "wb") as image_file:
                image_file.write(image_bytes)

            try:
                # 件进行 OCR
                text = pytesseract.image_to_string(image, lang= "chi_tra+eng")
            except Exception as e:
                print(f"Error processing image on slide {slide_number}: {e}")
            # 将图片和对应的文字一起存储
            images_with_text.append({'image_filename': image_filename, 'text': text})
    return images_with_text

def process_pptx_documents(pptx_dir):
    pptx_files = [os.path.join(pptx_dir, f) for f in os.listdir(pptx_dir) if f.endswith('.pptx')]

    try:
        with st.spinner("製作向量資料庫..."):
            all_docs = []
            for pptx_path in pptx_files:
                pptx_name = os.path.basename(pptx_path)
                presentation = Presentation(pptx_path)

                # 遍历每一张幻灯片
                for slide_number, slide in enumerate(presentation.slides, start=1):
                    # 提取幻灯片中的文本数据
                    text = extract_text_from_slide(slide)

                    # 提取幻灯片中的图片
                    images = extract_images_from_slide(slide, slide_number,pptx_name)

                    # 提取幻灯片中的表格数据
                    tables = extract_tables_from_slide(slide)

                    # 创建单张幻灯片的元数据
                    slide_metadata = {
                        'slide_number': slide_number,
                        'text': text,
                        'images': images,
                        'tables': tables
                    }

                    # 将每张幻灯片的信息附加到文档列表中
                    all_docs.append({
                        'source': pptx_name,
                        'slide_metadata': slide_metadata
                    })
            print(all_docs)
            # 在此处将 all_docs 传递给你的文档分割和检索器创建函数
            # docs = split_documents(all_docs)
            # st.session_state.retriever = create_retriever(docs)
            # st.session_state.retriever = parent_doc_retrevier(docs)
            # st.session_state.retriever = create_ensemble_retriever(all_docs)

        # st.success("資料庫製作完成")
        print("資料庫製作完成")
    except Exception as e:
        print(e)
        # st.error(f"處理文件時出現錯誤: {e}")




# 示例用法
pptx_dir = "./pptdata"
process_pptx_documents(pptx_dir)
